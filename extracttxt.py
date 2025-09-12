from pptx import Presentation
import sys
import os

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extracttxt.py <pptx_file>")
        sys.exit(1)
    pptx_file = sys.argv[1]
    if not os.path.isfile(pptx_file):
        print(f"File not found: {pptx_file}")
        sys.exit(1)
    extracted = extract_text_from_pptx(pptx_file)
    txt_file = os.path.splitext(pptx_file)[0] + ".txt"
    with open(txt_file, "w", encoding="utf-8") as f:
        f.write(extracted)
    print(f"Extracted text saved to {txt_file}")
